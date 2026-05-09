"""
tools/doc_parser.py

Multi-format document text extractor.
Supports: PDF, DOCX (Word), PPTX (PowerPoint), TXT, MD

Returns a list of source dicts compatible with the ResearchState.
"""
from __future__ import annotations

import os
import io
from typing import Any


def extract_text_from_file(file_bytes: bytes, filename: str) -> list[dict]:
    """
    Extract text from uploaded file bytes and return as a list of source dicts.

    Parameters
    ----------
    file_bytes : raw bytes of the uploaded file
    filename   : original filename (used to detect format)

    Returns
    -------
    List of source dicts with keys: url, title, snippet, source_type
    """
    ext = os.path.splitext(filename)[1].lower()
    sources = []

    try:
        if ext == ".pdf":
            sources = _parse_pdf(file_bytes, filename)
        elif ext in (".docx", ".doc"):
            sources = _parse_docx(file_bytes, filename)
        elif ext in (".pptx", ".ppt"):
            sources = _parse_pptx(file_bytes, filename)
        elif ext in (".txt", ".md", ".csv"):
            sources = _parse_text(file_bytes, filename)
        else:
            sources = [{
                "url": f"upload://{filename}",
                "title": filename,
                "snippet": f"[Unsupported format: {ext}. Supported: PDF, DOCX, PPTX, TXT, MD]",
                "source_type": "upload",
                "retrieved_at": _now(),
            }]
    except Exception as e:
        sources = [{
            "url": f"upload://{filename}",
            "title": filename,
            "snippet": f"[Error parsing file: {e}]",
            "source_type": "upload",
            "retrieved_at": _now(),
        }]

    return sources


# ── Per-format parsers ─────────────────────────────────────────────────────

def _parse_pdf(file_bytes: bytes, filename: str) -> list[dict]:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    sources = []
    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            sources.append({
                "url": f"upload://{filename}#page={i + 1}",
                "title": f"{filename} — Page {i + 1}",
                "snippet": text[:1200],
                "source_type": "pdf",
                "retrieved_at": _now(),
            })
    return sources or [_empty(filename, "PDF had no extractable text.")]


def _parse_docx(file_bytes: bytes, filename: str) -> list[dict]:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    chunks = []
    current = []
    for para in doc.paragraphs:
        txt = para.text.strip()
        if not txt:
            continue
        # Split into ~800-char chunks by heading
        if para.style.name.startswith("Heading") and current:
            chunks.append(" ".join(current))
            current = []
        current.append(txt)
    if current:
        chunks.append(" ".join(current))

    return [
        {
            "url": f"upload://{filename}#section={i + 1}",
            "title": f"{filename} — Section {i + 1}",
            "snippet": chunk[:1200],
            "source_type": "docx",
            "retrieved_at": _now(),
        }
        for i, chunk in enumerate(chunks)
    ] or [_empty(filename, "DOCX had no extractable text.")]


def _parse_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    sources = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        combined = " | ".join(texts)
        if combined:
            sources.append({
                "url": f"upload://{filename}#slide={i + 1}",
                "title": f"{filename} — Slide {i + 1}",
                "snippet": combined[:1200],
                "source_type": "pptx",
                "retrieved_at": _now(),
            })
    return sources or [_empty(filename, "PPTX had no extractable text.")]


def _parse_text(file_bytes: bytes, filename: str) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace").strip()
    # Split into 1200-char chunks
    chunks = [text[i:i + 1200] for i in range(0, len(text), 1200)]
    return [
        {
            "url": f"upload://{filename}#chunk={i + 1}",
            "title": f"{filename} — Part {i + 1}",
            "snippet": chunk,
            "source_type": "text",
            "retrieved_at": _now(),
        }
        for i, chunk in enumerate(chunks)
    ] or [_empty(filename, "File was empty.")]


# ── Helpers ────────────────────────────────────────────────────────────────

def _now() -> str:
    from datetime import datetime, timezone
    return datetime.now(timezone.utc).isoformat()


def _empty(filename: str, reason: str) -> dict:
    return {
        "url": f"upload://{filename}",
        "title": filename,
        "snippet": f"[{reason}]",
        "source_type": "upload",
        "retrieved_at": _now(),
    }
